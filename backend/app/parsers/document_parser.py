"""
Document Parser Module for Medical Handwriting OCR.

Handles parsing of printed medical documents in various formats including
PDF, DOCX, PPTX, and HTML.  Uses Marker for PDF conversion and Surya
for table recognition, with graceful fallbacks when libraries are
unavailable.

Supports Arabic (RTL) text extraction and proper bidirectional handling.
"""

from __future__ import annotations

import logging
import os
import tempfile
import uuid
from pathlib import Path
from typing import List, Optional

from pydantic import BaseModel, Field

from app.config import settings

logger = logging.getLogger(__name__)


# =============================================================================
# Pydantic Models
# =============================================================================


class ImageContent(BaseModel):
    """Represents an extracted image from a document page."""

    page_number: int = Field(..., description="1-based page number where image was found")
    image_path: str = Field(..., description="File path to the extracted image")
    bbox: Optional[dict] = Field(
        None,
        description="Bounding box {x1, y1, x2, y2} relative to the page",
    )
    caption: Optional[str] = Field(None, description="Associated caption text if available")
    width: Optional[int] = Field(None, description="Image width in pixels")
    height: Optional[int] = Field(None, description="Image height in pixels")
    mime_type: Optional[str] = Field(None, description="Detected MIME type")


class TableContent(BaseModel):
    """Represents a detected table within a document page."""

    page_number: int = Field(..., description="1-based page number")
    table_index: int = Field(..., description="0-based index on the page")
    headers: List[str] = Field(default_factory=list, description="Column header texts")
    rows: List[List[str]] = Field(
        default_factory=list,
        description="Table body rows (list of cell strings)",
    )
    confidence: float = Field(
        default=0.0,
        ge=0.0,
        le=1.0,
        description="Table extraction confidence score",
    )
    bbox: Optional[dict] = Field(
        None,
        description="Bounding box {x1, y1, x2, y2} on the page",
    )
    row_count: int = Field(0, description="Number of data rows")
    col_count: int = Field(0, description="Number of columns")


class PageContent(BaseModel):
    """Represents the extracted content of a single document page."""

    page_number: int = Field(..., description="1-based page number")
    text: str = Field(default="", description="Full extracted text (may include Arabic)")
    images: List[ImageContent] = Field(default_factory=list)
    tables: List[TableContent] = Field(default_factory=list)
    word_count: int = Field(0, description="Approximate word count for the page")
    has_arabic: bool = Field(False, description="Whether Arabic characters were detected")
    has_latin: bool = Field(False, description="Whether Latin characters were detected")


class DocumentParseResult(BaseModel):
    """Aggregated result of parsing an entire document."""

    document_id: str = Field(
        default_factory=lambda: str(uuid.uuid4()),
        description="Unique identifier for this parse operation",
    )
    file_name: str = Field("", description="Original file name")
    file_type: str = Field("", description="File extension (pdf, docx, pptx, html)")
    page_count: int = Field(0, description="Total number of pages")
    pages: List[PageContent] = Field(default_factory=list)
    full_text: str = Field("", description="Concatenated text from all pages")
    total_tables: int = Field(0, description="Total tables extracted")
    total_images: int = Field(0, description="Total images extracted")
    has_arabic: bool = Field(False, description="Whether any page contains Arabic text")
    processing_time_ms: float = Field(0.0, description="Total processing time in ms")
    warnings: List[str] = Field(default_factory=list, description="Non-fatal warnings")


# =============================================================================
# Helper utilities
# =============================================================================


def _has_arabic(text: str) -> bool:
    """Return True if *text* contains any Arabic character."""
    return any("\u0600" <= c <= "\u06FF" for c in text)


def _has_latin(text: str) -> bool:
    """Return True if *text* contains any Latin alphabetic character."""
    return any(c.isascii() and c.isalpha() for c in text)


def _add_rtl_markers(text: str) -> str:
    """
    Wrap detected Arabic segments with Unicode BiDi markers so that
    downstream consumers render RTL text correctly.
    """
    if not text:
        return text
    lines: list[str] = []
    for line in text.split("\n"):
        if _has_arabic(line):
            lines.append(f"\u202B{line}\u202C")
        else:
            lines.append(line)
    return "\n".join(lines)


# =============================================================================
# DocumentParser
# =============================================================================


class DocumentParser:
    """
    Multi-format document parser for printed medical documents.

    Supported formats:
    * **PDF** – parsed via ``marker-pdf`` with Surya OCR fallback
    * **DOCX** – parsed via ``python-docx``
    * **PPTX** – parsed via ``python-pptx``
    * **HTML** – parsed via ``beautifulsoup4``

    Arabic text is handled with proper RTL (right-to-left) markers.
    """

    def __init__(self) -> None:
        self._marker_available: Optional[bool] = None
        self._docx_available: Optional[bool] = None
        self._pptx_available: Optional[bool] = None
        self._bs4_available: Optional[bool] = None
        self._pdf2image_available: Optional[bool] = None
        self._output_dir: str = str(Path(settings.UPLOAD_DIR) / "parsed")
        os.makedirs(self._output_dir, exist_ok=True)
        logger.info("DocumentParser initialized. Output dir: %s", self._output_dir)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def parse_document(
        self,
        file_path: str,
        file_type: Optional[str] = None,
    ) -> DocumentParseResult:
        """
        Parse a document file and return structured content.

        Parameters
        ----------
        file_path : str
            Absolute or relative path to the document file.
        file_type : str, optional
            Override file type (``pdf``, ``docx``, ``pptx``, ``html``).
            If ``None`` the type is inferred from the file extension.

        Returns
        -------
        DocumentParseResult
            Structured parse result with pages, tables, images, and text.
        """
        import time

        start = time.perf_counter()
        file_type = (file_type or "").lower().lstrip(".")
        if not file_type:
            file_type = Path(file_path).suffix.lstrip(".").lower()

        file_name = os.path.basename(file_path)
        logger.info("Parsing document: %s (type=%s)", file_name, file_type)

        result = DocumentParseResult(
            file_name=file_name,
            file_type=file_type,
        )

        try:
            if file_type == "pdf":
                result = self._parse_pdf(file_path, result)
            elif file_type == "docx":
                result = self._parse_docx(file_path, result)
            elif file_type in ("pptx", "ppt"):
                result = self._parse_pptx(file_path, result)
            elif file_type in ("html", "htm"):
                result = self._parse_html(file_path, result)
            else:
                result.warnings.append(f"Unsupported file type: {file_type}. Attempting PDF fallback.")
                result = self._parse_pdf(file_path, result)

        except Exception as exc:
            logger.error("Failed to parse %s: %s", file_name, exc, exc_info=True)
            result.warnings.append(f"Parse error: {exc}")

        # Compute aggregated stats
        result.page_count = len(result.pages)
        result.full_text = "\n\n".join(p.text for p in result.pages)
        result.total_tables = sum(len(p.tables) for p in result.pages)
        result.total_images = sum(len(p.images) for p in result.pages)
        result.has_arabic = any(p.has_arabic for p in result.pages)
        result.processing_time_ms = (time.perf_counter() - start) * 1000

        logger.info(
            "Document parsed: %d pages, %d tables, %d images, %.1fms, arabic=%s",
            result.page_count,
            result.total_tables,
            result.total_images,
            result.processing_time_ms,
            result.has_arabic,
        )
        return result

    def extract_text_by_page(self, file_path: str) -> List[str]:
        """
        Quick helper: return a list of page-level text strings.

        Parameters
        ----------
        file_path : str
            Path to the document.

        Returns
        -------
        list[str]
            One string per page, in page order.
        """
        result = self.parse_document(file_path)
        return [page.text for page in result.pages]

    def extract_images(self, file_path: str) -> List[str]:
        """
        Quick helper: extract and save all images from a document.

        Parameters
        ----------
        file_path : str
            Path to the document.

        Returns
        -------
        list[str]
            List of file paths to the extracted images.
        """
        result = self.parse_document(file_path)
        paths: List[str] = []
        for page in result.pages:
            for img in page.images:
                paths.append(img.image_path)
        return paths

    def extract_tables(self, file_path: str) -> List[dict]:
        """
        Quick helper: extract all tables from a document.

        Parameters
        ----------
        file_path : str
            Path to the document.

        Returns
        -------
        list[dict]
            List of table dicts with ``headers``, ``rows``, ``page_number``.
        """
        result = self.parse_document(file_path)
        tables: List[dict] = []
        for page in result.pages:
            for tbl in page.tables:
                tables.append({
                    "page_number": tbl.page_number,
                    "headers": tbl.headers,
                    "rows": tbl.rows,
                    "row_count": tbl.row_count,
                    "col_count": tbl.col_count,
                    "confidence": tbl.confidence,
                })
        return tables

    # ------------------------------------------------------------------
    # PDF parsing (Marker + fallbacks)
    # ------------------------------------------------------------------

    def _parse_pdf(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Parse a PDF document using Marker with Surya fallback."""
        if self._check_marker():
            return self._parse_pdf_marker(file_path, result)

        # Fallback: use pdf2image + basic text extraction
        result.warnings.append(
            "marker-pdf not available; falling back to basic PDF text extraction"
        )
        return self._parse_pdf_fallback(file_path, result)

    def _check_marker(self) -> bool:
        """Lazily check whether marker-pdf is importable."""
        if self._marker_available is None:
            try:
                import marker  # noqa: F401
                self._marker_available = True
                logger.info("marker-pdf is available")
            except ImportError:
                self._marker_available = False
                logger.warning("marker-pdf is not installed")
        return self._marker_available

    def _parse_pdf_marker(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Parse PDF using marker-pdf library."""
        try:
            from marker.convert import convert_single_pdf
            from marker.models import load_all_models
        except ImportError:
            result.warnings.append("marker-pdf import failed despite availability check")
            return self._parse_pdf_fallback(file_path, result)

        try:
            load_all_models()
            full_text, _, images = convert_single_pdf(file_path)

            page = PageContent(
                page_number=1,
                text=_add_rtl_markers(full_text or ""),
                has_arabic=_has_arabic(full_text or ""),
                has_latin=_has_latin(full_text or ""),
                word_count=len((full_text or "").split()),
            )
            result.pages.append(page)

            if images:
                for idx, img_data in enumerate(images):
                    img_path = os.path.join(
                        self._output_dir,
                        f"{uuid.uuid4().hex}_img{idx}.png",
                    )
                    with open(img_path, "wb") as f:
                        f.write(img_data)
                    page.images.append(
                        ImageContent(
                            page_number=1,
                            image_path=img_path,
                            image_index=idx,
                        )
                    )

            logger.info("marker-pdf parsed successfully: %d chars", len(full_text or ""))
            return result

        except Exception as exc:
            logger.warning("marker-pdf parsing failed: %s", exc)
            result.warnings.append(f"marker-pdf error: {exc}")
            return self._parse_pdf_fallback(file_path, result)

    def _parse_pdf_fallback(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Fallback PDF parser using PyMuPDF (fitz) or pdfplumber."""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)

            for page_idx in range(len(doc)):
                page = doc.load_page(page_idx)
                text = page.get_text("text") or ""

                page_content = PageContent(
                    page_number=page_idx + 1,
                    text=_add_rtl_markers(text),
                    has_arabic=_has_arabic(text),
                    has_latin=_has_latin(text),
                    word_count=len(text.split()),
                )

                # Extract embedded images
                image_list = page.get_images(full=True)
                for img_idx, img_info in enumerate(image_list):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        if base_image:
                            img_bytes = base_image["image"]
                            img_ext = base_image.get("ext", "png")
                            img_path = os.path.join(
                                self._output_dir,
                                f"{uuid.uuid4().hex}_p{page_idx + 1}_i{img_idx}.{img_ext}",
                            )
                            with open(img_path, "wb") as f:
                                f.write(img_bytes)
                            page_content.images.append(
                                ImageContent(
                                    page_number=page_idx + 1,
                                    image_path=img_path,
                                    width=base_image.get("width"),
                                    height=base_image.get("height"),
                                    mime_type=base_image.get("mime"),
                                )
                            )
                    except Exception as img_exc:
                        logger.debug("Failed to extract image %d from page %d: %s", img_idx, page_idx + 1, img_exc)

                result.pages.append(page_content)

            doc.close()
            logger.info("PyMuPDF fallback parsed %d pages", len(result.pages))
            return result

        except ImportError:
            pass

        # Second fallback: pdfplumber
        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    page_content = PageContent(
                        page_number=page_idx + 1,
                        text=_add_rtl_markers(text),
                        has_arabic=_has_arabic(text),
                        has_latin=_has_latin(text),
                        word_count=len(text.split()),
                    )
                    result.pages.append(page_content)

            logger.info("pdfplumber fallback parsed %d pages", len(result.pages))
            return result

        except ImportError:
            result.warnings.append(
                "Neither PyMuPDF nor pdfplumber available for PDF parsing"
            )
            return result

    # ------------------------------------------------------------------
    # DOCX parsing
    # ------------------------------------------------------------------

    def _parse_docx(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Parse a DOCX file using python-docx."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            result.warnings.append("python-docx is not installed")
            return result

        try:
            doc = DocxDocument(file_path)
            page_num = 1
            buffer_lines: List[str] = []

            for para in doc.paragraphs:
                text = para.text or ""
                if text.strip():
                    buffer_lines.append(text)

                # Treat page breaks as page separators
                if any(
                    run._element.xml.find("w:br") != -1
                    and 'w:type="page"' in run._element.xml
                    for run in para.runs
                ) if para.runs else False:
                    self._flush_page(buffer_lines, page_num, result)
                    page_num += 1
                    buffer_lines = []

            # Flush remaining text as the last page
            if buffer_lines:
                self._flush_page(buffer_lines, page_num, result)

            # Extract images from DOCX
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    try:
                        image_data = rel.target_part.blob
                        ext = rel.target_part.content_type.split("/")[-1]
                        ext = {"jpeg": "jpg", "png": "png", "gif": "gif", "tiff": "tiff"}.get(ext, "png")
                        img_path = os.path.join(
                            self._output_dir,
                            f"{uuid.uuid4().hex}.{ext}",
                        )
                        with open(img_path, "wb") as f:
                            f.write(image_data)
                        if result.pages:
                            result.pages[0].images.append(
                                ImageContent(
                                    page_number=1,
                                    image_path=img_path,
                                    mime_type=rel.target_part.content_type,
                                )
                            )
                    except Exception as img_exc:
                        logger.debug("Failed to extract DOCX image: %s", img_exc)

            logger.info("DOCX parsed: %d pages", len(result.pages))
            return result

        except Exception as exc:
            logger.error("DOCX parsing failed: %s", exc, exc_info=True)
            result.warnings.append(f"DOCX parse error: {exc}")
            return result

    def _flush_page(
        self, lines: List[str], page_num: int, result: DocumentParseResult
    ) -> None:
        """Helper to flush accumulated text lines into a PageContent."""
        text = "\n".join(lines)
        page = PageContent(
            page_number=page_num,
            text=_add_rtl_markers(text),
            has_arabic=_has_arabic(text),
            has_latin=_has_latin(text),
            word_count=len(text.split()),
        )
        result.pages.append(page)

    # ------------------------------------------------------------------
    # PPTX parsing
    # ------------------------------------------------------------------

    def _parse_pptx(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Parse a PPTX file using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            result.warnings.append("python-pptx is not installed")
            return result

        try:
            prs = Presentation(file_path)

            for slide_idx, slide in enumerate(prs.slides):
                page_num = slide_idx + 1
                slide_texts: List[str] = []

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text or ""
                            if para_text.strip():
                                slide_texts.append(para_text)

                    # Extract images from shapes
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            image_blob = shape.image.blob
                            ext = shape.image.content_type.split("/")[-1]
                            ext = {"jpeg": "jpg", "png": "png", "gif": "gif"}.get(ext, "png")
                            img_path = os.path.join(
                                self._output_dir,
                                f"{uuid.uuid4().hex}_slide{page_num}.{ext}",
                            )
                            with open(img_path, "wb") as f:
                                f.write(image_blob)
                            if result.pages:
                                result.pages[-1].images.append(
                                    ImageContent(
                                        page_number=page_num,
                                        image_path=img_path,
                                        mime_type=shape.image.content_type,
                                    )
                                )
                        except Exception as img_exc:
                            logger.debug("Failed to extract PPTX image from slide %d: %s", page_num, img_exc)

                text = "\n".join(slide_texts)
                page = PageContent(
                    page_number=page_num,
                    text=_add_rtl_markers(text),
                    has_arabic=_has_arabic(text),
                    has_latin=_has_latin(text),
                    word_count=len(text.split()),
                )
                result.pages.append(page)

            logger.info("PPTX parsed: %d slides", len(result.pages))
            return result

        except Exception as exc:
            logger.error("PPTX parsing failed: %s", exc, exc_info=True)
            result.warnings.append(f"PPTX parse error: {exc}")
            return result

    # ------------------------------------------------------------------
    # HTML parsing
    # ------------------------------------------------------------------

    def _parse_html(
        self, file_path: str, result: DocumentParseResult
    ) -> DocumentParseResult:
        """Parse an HTML file using BeautifulSoup."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            result.warnings.append("beautifulsoup4 is not installed")
            return result

        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, "html.parser")

            # Remove script/style tags
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()

            # Extract text preserving some structure
            body = soup.find("body") or soup
            text = body.get_text(separator="\n", strip=True)
            text = _add_rtl_markers(text)

            page = PageContent(
                page_number=1,
                text=text,
                has_arabic=_has_arabic(text),
                has_latin=_has_latin(text),
                word_count=len(text.split()),
            )
            result.pages.append(page)

            # Extract images
            for img_tag in soup.find_all("img"):
                src = img_tag.get("src", "")
                alt = img_tag.get("alt", "")
                if src:
                    # Save as a reference; actual download is beyond scope here
                    page.images.append(
                        ImageContent(
                            page_number=1,
                            image_path=src,
                            caption=alt or None,
                        )
                    )

            # Try to detect tables in HTML
            for table_idx, table_tag in enumerate(soup.find_all("table")):
                headers: List[str] = []
                rows: List[List[str]] = []

                # First row as headers
                header_row = table_tag.find("tr")
                if header_row:
                    for th in header_row.find_all(["th", "td"]):
                        headers.append(th.get_text(strip=True))

                # Remaining rows
                for tr in table_tag.find_all("tr")[1:]:
                    row_cells = [td.get_text(strip=True) for td in tr.find_all(["th", "td"])]
                    if row_cells:
                        rows.append(row_cells)

                if headers or rows:
                    page.tables.append(
                        TableContent(
                            page_number=1,
                            table_index=table_idx,
                            headers=headers,
                            rows=rows,
                            row_count=len(rows),
                            col_count=max(len(headers), max((len(r) for r in rows), default=0)),
                            confidence=0.9 if rows else 0.5,
                        )
                    )

            logger.info("HTML parsed: %d chars, %d tables", len(text), len(page.tables))
            return result

        except Exception as exc:
            logger.error("HTML parsing failed: %s", exc, exc_info=True)
            result.warnings.append(f"HTML parse error: {exc}")
            return result


# =============================================================================
# Singleton instance
# =============================================================================

document_parser = DocumentParser()
